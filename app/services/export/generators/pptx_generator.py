"""
Comprehensive PowerPoint/PPTX generator for academic presentations.

This module provides a complete PPTX generation system with:
- Academic template system (IEEE, ACM, Nature, university themes)
- Custom branding support
- Advanced features (transitions, speaker notes, equations, images)
- Citation integration
- Layout management with responsive design
- High-quality image handling
"""

import io
import logging
import math
import re
import tempfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union
from urllib.parse import urlparse
from uuid import UUID

try:
    import requests
    from PIL import Image, ImageDraw, ImageFont
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls, qn
    from pptx.parts.image import Image as PptxImage
    from pptx.slide import Slide
    from pptx.text.text import TextFrame
    from pptx.util import Cm, Inches, Pt
except ImportError as e:
    raise ImportError(
        f"Required packages not installed: {e}. "
        "Please install: pip install python-pptx pillow requests"
    )

from app.core.logging import get_logger
from app.domain.schemas.generation import Citation, SlideContent

logger = get_logger(__name__)


class AcademicTemplate(Enum):
    """Academic presentation templates."""
    IEEE = "ieee"
    ACM = "acm"
    NATURE = "nature"
    SPRINGER = "springer"
    ELSEVIER = "elsevier"
    MIT = "mit"
    STANFORD = "stanford"
    HARVARD = "harvard"
    OXFORD = "oxford"
    CAMBRIDGE = "cambridge"
    CUSTOM = "custom"


class SlideLayout(Enum):
    """Slide layout types."""
    TITLE_SLIDE = "title_slide"
    TITLE_CONTENT = "title_content"
    TWO_CONTENT = "two_content"
    CONTENT_WITH_CAPTION = "content_with_caption"
    BLANK = "blank"
    SECTION_HEADER = "section_header"
    COMPARISON = "comparison"
    PICTURE_WITH_CAPTION = "picture_with_caption"


class TransitionType(Enum):
    """Slide transition types."""
    NONE = "none"
    FADE = "fade"
    PUSH = "push"
    WIPE = "wipe"
    SPLIT = "split"
    REVEAL = "reveal"
    RANDOM_BARS = "random_bars"
    SHAPE = "shape"
    UNCOVER = "uncover"
    COVER = "cover"


@dataclass
class ColorScheme:
    """Color scheme for presentations."""
    primary: str = "#1f4e79"  # Dark blue
    secondary: str = "#70ad47"  # Green
    accent: str = "#c5504b"  # Red
    text_primary: str = "#000000"  # Black
    text_secondary: str = "#404040"  # Dark gray
    background: str = "#ffffff"  # White
    background_alt: str = "#f2f2f2"  # Light gray


@dataclass
class Typography:
    """Typography settings."""
    title_font: str = "Calibri"
    body_font: str = "Calibri"
    code_font: str = "Consolas"
    title_size: int = 44
    subtitle_size: int = 32
    body_size: int = 24
    caption_size: int = 18
    footnote_size: int = 14


@dataclass
class BrandingConfig:
    """Branding configuration."""
    logo_path: Optional[str] = None
    logo_position: str = "top_right"  # top_left, top_right, bottom_left, bottom_right
    logo_size: Tuple[float, float] = (1.5, 0.75)  # inches
    university_name: Optional[str] = None
    department_name: Optional[str] = None
    show_slide_numbers: bool = True
    show_date: bool = True
    custom_footer: Optional[str] = None


@dataclass
class TemplateConfig:
    """Complete template configuration."""
    template: AcademicTemplate = AcademicTemplate.IEEE
    color_scheme: ColorScheme = field(default_factory=ColorScheme)
    typography: Typography = field(default_factory=Typography)
    branding: BrandingConfig = field(default_factory=BrandingConfig)
    slide_size: Tuple[float, float] = (16, 9)  # 16:9 aspect ratio in inches
    margins: Tuple[float, float, float, float] = (0.5, 0.5, 0.5, 0.5)  # left, top, right, bottom


class EquationRenderer:
    """Renders mathematical equations for presentations."""
    
    def __init__(self):
        self.logger = get_logger(self.__class__.__name__)
    
    def render_latex_to_image(self, latex_code: str, font_size: int = 24) -> Optional[io.BytesIO]:
        """
        Render LaTeX equation to image.
        
        Args:
            latex_code: LaTeX equation code
            font_size: Font size for rendering
            
        Returns:
            BytesIO object containing the rendered image
        """
        try:
            # For production, you might want to use matplotlib or a LaTeX service
            # This is a simplified version using PIL
            img = Image.new('RGBA', (400, 100), (255, 255, 255, 0))
            draw = ImageDraw.Draw(img)
            
            # Simple text rendering as fallback
            # In production, integrate with matplotlib's mathtext or external LaTeX service
            clean_text = re.sub(r'[{}\\]', '', latex_code)
            try:
                font = ImageFont.truetype("arial.ttf", font_size)
            except OSError:
                font = ImageFont.load_default()
            
            draw.text((10, 10), clean_text, fill=(0, 0, 0, 255), font=font)
            
            # Crop to content
            bbox = img.getbbox()
            if bbox:
                img = img.crop(bbox)
            
            buffer = io.BytesIO()
            img.save(buffer, format='PNG')
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            self.logger.error(f"Failed to render equation: {e}")
            return None
    
    def extract_equations_from_text(self, text: str) -> List[Tuple[str, str]]:
        """
        Extract LaTeX equations from text.
        
        Args:
            text: Text containing LaTeX equations
            
        Returns:
            List of (equation_text, latex_code) tuples
        """
        equations = []
        
        # Match inline equations: $...$
        inline_pattern = r'\$([^$]+)\$'
        for match in re.finditer(inline_pattern, text):
            equations.append((match.group(0), match.group(1)))
        
        # Match display equations: $$...$$
        display_pattern = r'\$\$([^$]+)\$\$'
        for match in re.finditer(display_pattern, text):
            equations.append((match.group(0), match.group(1)))
        
        # Match LaTeX environments: \begin{equation}...\end{equation}
        env_pattern = r'\\begin\{(equation|align|gather)\*?\}(.*?)\\end\{\1\*?\}'
        for match in re.finditer(env_pattern, text, re.DOTALL):
            equations.append((match.group(0), match.group(2).strip()))
        
        return equations


class ImageProcessor:
    """Processes and optimizes images for presentations."""
    
    def __init__(self):
        self.logger = get_logger(self.__class__.__name__)
        self.max_image_size = (1920, 1080)  # Max resolution
        self.supported_formats = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff'}
    
    def process_image(
        self, 
        image_path: Union[str, io.BytesIO], 
        target_size: Optional[Tuple[int, int]] = None,
        maintain_aspect: bool = True
    ) -> io.BytesIO:
        """
        Process and optimize image for presentation.
        
        Args:
            image_path: Path to image file or BytesIO object
            target_size: Target size (width, height)
            maintain_aspect: Whether to maintain aspect ratio
            
        Returns:
            Processed image as BytesIO
        """
        try:
            if isinstance(image_path, str):
                if image_path.startswith(('http://', 'https://')):
                    # Download image
                    response = requests.get(image_path, timeout=30)
                    response.raise_for_status()
                    img = Image.open(io.BytesIO(response.content))
                else:
                    img = Image.open(image_path)
            else:
                img = Image.open(image_path)
            
            # Convert to RGB if necessary
            if img.mode in ('RGBA', 'LA', 'P'):
                # Create white background
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1])
                img = background
            elif img.mode != 'RGB':
                img = img.convert('RGB')
            
            # Resize if needed
            if target_size:
                if maintain_aspect:
                    img.thumbnail(target_size, Image.Resampling.LANCZOS)
                else:
                    img = img.resize(target_size, Image.Resampling.LANCZOS)
            else:
                # Apply max size constraint
                img.thumbnail(self.max_image_size, Image.Resampling.LANCZOS)
            
            # Save to BytesIO
            buffer = io.BytesIO()
            img.save(buffer, format='JPEG', quality=90, optimize=True)
            buffer.seek(0)
            return buffer
            
        except Exception as e:
            self.logger.error(f"Failed to process image: {e}")
            # Return a placeholder image
            return self._create_placeholder_image(target_size or (800, 600))
    
    def _create_placeholder_image(self, size: Tuple[int, int]) -> io.BytesIO:
        """Create a placeholder image."""
        img = Image.new('RGB', size, (240, 240, 240))
        draw = ImageDraw.Draw(img)
        
        # Draw placeholder text
        text = "Image Placeholder"
        try:
            font = ImageFont.truetype("arial.ttf", size[1] // 20)
        except OSError:
            font = ImageFont.load_default()
        
        text_bbox = draw.textbbox((0, 0), text, font=font)
        text_width = text_bbox[2] - text_bbox[0]
        text_height = text_bbox[3] - text_bbox[1]
        
        x = (size[0] - text_width) // 2
        y = (size[1] - text_height) // 2
        draw.text((x, y), text, fill=(128, 128, 128), font=font)
        
        # Draw border
        draw.rectangle([0, 0, size[0]-1, size[1]-1], outline=(180, 180, 180), width=2)
        
        buffer = io.BytesIO()
        img.save(buffer, format='JPEG', quality=90)
        buffer.seek(0)
        return buffer


class CitationFormatter:
    """Formats citations according to academic styles."""
    
    def __init__(self):
        self.logger = get_logger(self.__class__.__name__)
    
    def format_citation(self, citation: Citation, style: str = 'ieee') -> str:
        """
        Format a single citation.
        
        Args:
            citation: Citation object
            style: Citation style (ieee, apa, mla, chicago)
            
        Returns:
            Formatted citation string
        """
        if style.lower() == 'ieee':
            return self._format_ieee(citation)
        elif style.lower() == 'apa':
            return self._format_apa(citation)
        elif style.lower() == 'mla':
            return self._format_mla(citation)
        elif style.lower() == 'chicago':
            return self._format_chicago(citation)
        else:
            return self._format_ieee(citation)  # Default to IEEE
    
    def _format_ieee(self, citation: Citation) -> str:
        """Format citation in IEEE style."""
        authors = self._format_authors_ieee(citation.authors)
        title = f'"{citation.title}"'
        
        parts = [authors, title]
        
        if citation.venue:
            parts.append(citation.venue)
        
        if citation.year:
            parts.append(str(citation.year))
        
        if citation.doi:
            parts.append(f"doi: {citation.doi}")
        
        return ', '.join(filter(None, parts)) + '.'
    
    def _format_apa(self, citation: Citation) -> str:
        """Format citation in APA style."""
        authors = self._format_authors_apa(citation.authors)
        year = f"({citation.year})" if citation.year else ""
        title = citation.title
        
        parts = [authors, year, title]
        
        if citation.venue:
            parts.append(f"*{citation.venue}*")
        
        if citation.doi:
            parts.append(f"https://doi.org/{citation.doi}")
        
        return '. '.join(filter(None, parts)) + '.'
    
    def _format_mla(self, citation: Citation) -> str:
        """Format citation in MLA style."""
        if citation.authors:
            author = citation.authors[0]
            if len(citation.authors) > 1:
                author += " et al."
        else:
            author = "Unknown Author"
        
        parts = [author, f'"{citation.title}"']
        
        if citation.venue:
            parts.append(f"*{citation.venue}*")
        
        if citation.year:
            parts.append(str(citation.year))
        
        return ', '.join(filter(None, parts)) + '.'
    
    def _format_chicago(self, citation: Citation) -> str:
        """Format citation in Chicago style."""
        authors = self._format_authors_chicago(citation.authors)
        title = f'"{citation.title}"'
        
        parts = [authors, title]
        
        if citation.venue:
            parts.append(citation.venue)
        
        if citation.year:
            parts.append(f"({citation.year})")
        
        if citation.doi:
            parts.append(f"doi:{citation.doi}")
        
        return '. '.join(filter(None, parts)) + '.'
    
    def _format_authors_ieee(self, authors: List[str]) -> str:
        """Format authors for IEEE style."""
        if not authors:
            return ""
        
        if len(authors) == 1:
            return authors[0]
        elif len(authors) <= 3:
            return ', '.join(authors[:-1]) + ' and ' + authors[-1]
        else:
            return authors[0] + ' et al.'
    
    def _format_authors_apa(self, authors: List[str]) -> str:
        """Format authors for APA style."""
        if not authors:
            return ""
        
        # Convert to Last, F. M. format
        formatted_authors = []
        for author in authors:
            if ',' in author:
                formatted_authors.append(author)
            else:
                parts = author.split()
                if len(parts) >= 2:
                    last = parts[-1]
                    first_initials = '. '.join([p[0] for p in parts[:-1]]) + '.'
                    formatted_authors.append(f"{last}, {first_initials}")
                else:
                    formatted_authors.append(author)
        
        if len(formatted_authors) == 1:
            return formatted_authors[0]
        elif len(formatted_authors) <= 7:
            return ', '.join(formatted_authors[:-1]) + ', & ' + formatted_authors[-1]
        else:
            return ', '.join(formatted_authors[:6]) + ', ... ' + formatted_authors[-1]
    
    def _format_authors_chicago(self, authors: List[str]) -> str:
        """Format authors for Chicago style."""
        if not authors:
            return ""
        
        if len(authors) == 1:
            return authors[0]
        elif len(authors) <= 3:
            return ', '.join(authors[:-1]) + ', and ' + authors[-1]
        else:
            return authors[0] + ' et al.'


class LayoutManager:
    """Manages slide layouts and content positioning."""
    
    def __init__(self, template_config: TemplateConfig):
        self.config = template_config
        self.logger = get_logger(self.__class__.__name__)
    
    def calculate_content_area(self, slide: Slide) -> Tuple[float, float, float, float]:
        """
        Calculate available content area considering margins and branding.
        
        Args:
            slide: PowerPoint slide
            
        Returns:
            (left, top, width, height) in inches
        """
        slide_width = self.config.slide_size[0]
        slide_height = self.config.slide_size[1]
        margins = self.config.margins
        
        left = margins[0]
        top = margins[1]
        width = slide_width - margins[0] - margins[2]
        height = slide_height - margins[1] - margins[3]
        
        # Account for logo space
        if self.config.branding.logo_path:
            logo_height = self.config.branding.logo_size[1]
            if self.config.branding.logo_position.startswith('top'):
                top += logo_height + 0.2  # 0.2" padding
                height -= logo_height + 0.2
            elif self.config.branding.logo_position.startswith('bottom'):
                height -= logo_height + 0.2
        
        # Account for slide numbers and footer
        if self.config.branding.show_slide_numbers or self.config.branding.custom_footer:
            height -= 0.5  # Reserve space for footer
        
        return left, top, width, height
    
    def fit_text_to_shape(self, text_frame: TextFrame, max_font_size: int = 32) -> None:
        """
        Automatically fit text to text frame.
        
        Args:
            text_frame: PowerPoint text frame
            max_font_size: Maximum font size to use
        """
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        
        # Set maximum font size
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size and run.font.size.pt > max_font_size:
                    run.font.size = Pt(max_font_size)
    
    def create_two_column_layout(
        self, 
        slide: Slide, 
        left_content: Dict[str, Any], 
        right_content: Dict[str, Any],
        title: str = ""
    ) -> None:
        """
        Create a two-column layout on the slide.
        
        Args:
            slide: PowerPoint slide
            left_content: Content for left column
            right_content: Content for right column
            title: Optional title
        """
        left, top, width, height = self.calculate_content_area(slide)
        
        title_height = 0
        if title:
            title_height = 1.0  # Reserve 1" for title
            top += title_height
            height -= title_height
            
            # Add title
            title_shape = slide.shapes.add_textbox(
                Inches(left), Inches(self.config.margins[1]), 
                Inches(width), Inches(title_height)
            )
            title_frame = title_shape.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(self.config.typography.title_size)
            title_frame.paragraphs[0].font.name = self.config.typography.title_font
        
        column_width = (width - 0.5) / 2  # 0.5" gap between columns
        
        # Left column
        if left_content:
            self._add_content_to_area(
                slide, left, top, column_width, height, left_content
            )
        
        # Right column
        if right_content:
            self._add_content_to_area(
                slide, left + column_width + 0.5, top, column_width, height, right_content
            )
    
    def _add_content_to_area(
        self, 
        slide: Slide, 
        x: float, y: float, width: float, height: float,
        content: Dict[str, Any]
    ) -> None:
        """Add content to a specific area of the slide."""
        content_type = content.get('type', 'text')
        
        if content_type == 'text':
            text_box = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            text_frame = text_box.text_frame
            text_frame.text = content.get('content', '')
            self.fit_text_to_shape(text_frame, self.config.typography.body_size)
            
        elif content_type == 'bullet_list':
            text_box = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            text_frame = text_box.text_frame
            text_frame.clear()
            
            for item in content.get('items', []):
                p = text_frame.paragraphs[0] if not text_frame.paragraphs[0].text else text_frame.add_paragraph()
                p.text = item
                p.level = 0
                p.font.size = Pt(self.config.typography.body_size)
            
        elif content_type == 'image':
            image_path = content.get('path') or content.get('url')
            if image_path:
                try:
                    processor = ImageProcessor()
                    image_buffer = processor.process_image(
                        image_path, 
                        target_size=(int(width * 72), int(height * 72))  # Convert to pixels
                    )
                    slide.shapes.add_picture(
                        image_buffer, Inches(x), Inches(y), 
                        width=Inches(width), height=Inches(height)
                    )
                except Exception as e:
                    self.logger.error(f"Failed to add image: {e}")
                    # Add placeholder text
                    text_box = slide.shapes.add_textbox(
                        Inches(x), Inches(y), Inches(width), Inches(height)
                    )
                    text_box.text_frame.text = "[Image placeholder]"


class PPTXGenerator:
    """
    Comprehensive PowerPoint/PPTX generator for academic presentations.
    
    Features:
    - Academic template system with predefined themes
    - Custom branding support with logos and color schemes
    - Advanced slide layouts and responsive design
    - Equation rendering and integration
    - High-quality image processing and optimization
    - Speaker notes and citation management
    - Slide transitions and animations
    - Export optimization and error handling
    """
    
    def __init__(self, template_config: Optional[TemplateConfig] = None):
        """
        Initialize the PPTX generator.
        
        Args:
            template_config: Template configuration. Uses default if None.
        """
        self.config = template_config or TemplateConfig()
        self.logger = get_logger(self.__class__.__name__)
        self.equation_renderer = EquationRenderer()
        self.image_processor = ImageProcessor()
        self.citation_formatter = CitationFormatter()
        self.layout_manager = LayoutManager(self.config)
        
        # Template definitions
        self.templates = self._load_template_definitions()
        
        # Initialize presentation
        self.presentation: Optional[Presentation] = None
    
    def _load_template_definitions(self) -> Dict[AcademicTemplate, Dict[str, Any]]:
        """Load predefined academic template configurations."""
        templates = {
            AcademicTemplate.IEEE: {
                'color_scheme': ColorScheme(
                    primary='#003f7f',
                    secondary='#0066cc',
                    accent='#ff6600',
                    text_primary='#000000',
                    text_secondary='#333333',
                    background='#ffffff',
                    background_alt='#f8f9fa'
                ),
                'typography': Typography(
                    title_font='Arial',
                    body_font='Arial',
                    title_size=44,
                    body_size=24
                )
            },
            AcademicTemplate.ACM: {
                'color_scheme': ColorScheme(
                    primary='#0085ca',
                    secondary='#005580',
                    accent='#ff6900',
                    text_primary='#000000',
                    text_secondary='#333333',
                    background='#ffffff',
                    background_alt='#f5f5f5'
                ),
                'typography': Typography(
                    title_font='Helvetica',
                    body_font='Helvetica',
                    title_size=42,
                    body_size=22
                )
            },
            AcademicTemplate.NATURE: {
                'color_scheme': ColorScheme(
                    primary='#006633',
                    secondary='#009966',
                    accent='#cc3300',
                    text_primary='#000000',
                    text_secondary='#444444',
                    background='#ffffff',
                    background_alt='#f0f8f0'
                ),
                'typography': Typography(
                    title_font='Times New Roman',
                    body_font='Times New Roman',
                    title_size=40,
                    body_size=20
                )
            },
            AcademicTemplate.MIT: {
                'color_scheme': ColorScheme(
                    primary='#8c1515',
                    secondary='#b83a4b',
                    accent='#009639',
                    text_primary='#000000',
                    text_secondary='#333333',
                    background='#ffffff',
                    background_alt='#f4f4f4'
                ),
                'typography': Typography(
                    title_font='Source Sans Pro',
                    body_font='Source Sans Pro',
                    title_size=46,
                    body_size=26
                )
            },
            AcademicTemplate.STANFORD: {
                'color_scheme': ColorScheme(
                    primary='#8c1515',
                    secondary='#b83a4b',
                    accent='#009b77',
                    text_primary='#2e2d29',
                    text_secondary='#544948',
                    background='#ffffff',
                    background_alt='#f4f4f4'
                ),
                'typography': Typography(
                    title_font='Source Sans Pro',
                    body_font='Source Sans Pro',
                    title_size=44,
                    body_size=24
                )
            }
        }
        return templates
    
    def create_presentation(
        self, 
        slides: List[SlideContent], 
        citations: Optional[List[Citation]] = None,
        metadata: Optional[Dict[str, Any]] = None
    ) -> Presentation:
        """
        Create a complete PowerPoint presentation.
        
        Args:
            slides: List of slide content
            citations: List of citations to include
            metadata: Presentation metadata
            
        Returns:
            PowerPoint presentation object
        """
        try:
            self.logger.info(f"Creating presentation with {len(slides)} slides")
            
            # Initialize presentation
            self.presentation = Presentation()
            self._setup_slide_size()
            self._apply_template()
            
            # Clear default slides
            xml_slides = self.presentation.slides._sldIdLst
            for i in range(len(xml_slides)):
                self.presentation.slides._sldIdLst.remove(xml_slides[0])
            
            # Add slides
            for i, slide_content in enumerate(slides):
                try:
                    self._add_slide(slide_content, i + 1)
                except Exception as e:
                    self.logger.error(f"Error creating slide {i + 1}: {e}")
                    # Add error slide as fallback
                    self._add_error_slide(f"Error in slide {i + 1}: {str(e)}")
            
            # Add references slide if citations provided
            if citations:
                self._add_references_slide(citations)
            
            # Apply global formatting
            self._apply_global_formatting()
            
            # Add metadata
            if metadata:
                self._add_metadata(metadata)
            
            self.logger.info("Presentation created successfully")
            return self.presentation
            
        except Exception as e:
            self.logger.error(f"Failed to create presentation: {e}")
            raise
    
    def save_presentation(self, output_path: Union[str, io.BytesIO]) -> None:
        """
        Save the presentation to file or buffer.
        
        Args:
            output_path: Output file path or BytesIO buffer
        """
        if not self.presentation:
            raise ValueError("No presentation to save. Call create_presentation first.")
        
        try:
            if isinstance(output_path, str):
                self.presentation.save(output_path)
                self.logger.info(f"Presentation saved to {output_path}")
            else:
                self.presentation.save(output_path)
                self.logger.info("Presentation saved to buffer")
                
        except Exception as e:
            self.logger.error(f"Failed to save presentation: {e}")
            raise
    
    def _setup_slide_size(self) -> None:
        """Set up slide dimensions."""
        if self.presentation:
            # Set slide size (16:9 aspect ratio)
            slide_width = int(self.config.slide_size[0] * 914400)  # Convert inches to EMUs
            slide_height = int(self.config.slide_size[1] * 914400)
            
            self.presentation.slide_width = slide_width
            self.presentation.slide_height = slide_height
    
    def _apply_template(self) -> None:
        """Apply the selected academic template."""
        if self.config.template in self.templates:
            template_def = self.templates[self.config.template]
            
            # Apply color scheme
            if 'color_scheme' in template_def:
                self.config.color_scheme = template_def['color_scheme']
            
            # Apply typography
            if 'typography' in template_def:
                self.config.typography = template_def['typography']
    
    def _add_slide(self, slide_content: SlideContent, slide_number: int) -> Slide:
        """
        Add a slide to the presentation.
        
        Args:
            slide_content: Content for the slide
            slide_number: Slide number (1-based)
            
        Returns:
            Created slide
        """
        if not self.presentation:
            raise ValueError("Presentation not initialized")
        
        # Determine layout
        layout_type = self._determine_layout(slide_content)
        
        # Add slide with blank layout
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add content based on layout
        if layout_type == SlideLayout.TITLE_SLIDE:
            self._create_title_slide(slide, slide_content)
        elif layout_type == SlideLayout.SECTION_HEADER:
            self._create_section_header_slide(slide, slide_content)
        elif layout_type == SlideLayout.TWO_CONTENT:
            self._create_two_content_slide(slide, slide_content)
        else:
            self._create_content_slide(slide, slide_content)
        
        # Add branding elements
        self._add_branding_elements(slide, slide_number)
        
        # Add speaker notes
        if slide_content.speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_content.speaker_notes
        
        # Apply transitions
        self._apply_slide_transition(slide)
        
        return slide
    
    def _determine_layout(self, slide_content: SlideContent) -> SlideLayout:
        """Determine the best layout for slide content."""
        # Title slide detection
        if (slide_content.metadata.get('slide_type') == 'title' or 
            any(keyword in (slide_content.title or '').lower() 
                for keyword in ['title', 'presentation', 'introduction'])):
            return SlideLayout.TITLE_SLIDE
        
        # Section header detection
        if (slide_content.metadata.get('slide_type') == 'section' or
            len(slide_content.body) == 0):
            return SlideLayout.SECTION_HEADER
        
        # Two content detection
        if len(slide_content.body) == 2 and all(
            item.get('type') in ['text', 'bullet_list', 'image'] 
            for item in slide_content.body
        ):
            return SlideLayout.TWO_CONTENT
        
        return SlideLayout.TITLE_CONTENT
    
    def _create_title_slide(self, slide: Slide, content: SlideContent) -> None:
        """Create a title slide."""
        left, top, width, height = self.layout_manager.calculate_content_area(slide)
        
        # Title
        if content.title:
            title_height = height * 0.3
            title_shape = slide.shapes.add_textbox(
                Inches(left), Inches(top + height * 0.2), 
                Inches(width), Inches(title_height)
            )
            title_frame = title_shape.text_frame
            title_frame.text = content.title
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title_frame.paragraphs[0].font.size = Pt(self.config.typography.title_size)
            title_frame.paragraphs[0].font.name = self.config.typography.title_font
            title_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                self.config.color_scheme.primary.lstrip('#')
            )
        
        # Subtitle
        if content.subtitle:
            subtitle_top = top + height * 0.55
            subtitle_height = height * 0.2
            subtitle_shape = slide.shapes.add_textbox(
                Inches(left), Inches(subtitle_top), 
                Inches(width), Inches(subtitle_height)
            )
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.text = content.subtitle
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle_frame.paragraphs[0].font.size = Pt(self.config.typography.subtitle_size)
            subtitle_frame.paragraphs[0].font.name = self.config.typography.body_font
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                self.config.color_scheme.text_secondary.lstrip('#')
            )
        
        # Additional content (author, institution, etc.)
        if content.body:
            content_top = top + height * 0.8
            content_height = height * 0.15
            
            for item in content.body:
                if item.get('type') == 'text':
                    text_shape = slide.shapes.add_textbox(
                        Inches(left), Inches(content_top), 
                        Inches(width), Inches(content_height)
                    )
                    text_frame = text_shape.text_frame
                    text_frame.text = item.get('content', '')
                    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    text_frame.paragraphs[0].font.size = Pt(self.config.typography.body_size)
                    content_top += 0.5
    
    def _create_section_header_slide(self, slide: Slide, content: SlideContent) -> None:
        """Create a section header slide."""
        left, top, width, height = self.layout_manager.calculate_content_area(slide)
        
        # Add colored background shape
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(height * 0.3), 
            Inches(self.config.slide_size[0]), Inches(height * 0.4)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor.from_string(
            self.config.color_scheme.primary.lstrip('#')
        )
        bg_shape.line.fill.background()
        
        # Section title
        if content.title:
            title_shape = slide.shapes.add_textbox(
                Inches(left), Inches(height * 0.4), 
                Inches(width), Inches(height * 0.3)
            )
            title_frame = title_shape.text_frame
            title_frame.text = content.title
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title_frame.paragraphs[0].font.size = Pt(self.config.typography.title_size + 4)
            title_frame.paragraphs[0].font.name = self.config.typography.title_font
            title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
            title_frame.paragraphs[0].font.bold = True
    
    def _create_two_content_slide(self, slide: Slide, content: SlideContent) -> None:
        """Create a slide with two content areas."""
        if len(content.body) >= 2:
            self.layout_manager.create_two_column_layout(
                slide, content.body[0], content.body[1], content.title
            )
    
    def _create_content_slide(self, slide: Slide, content: SlideContent) -> None:
        """Create a standard content slide."""
        left, top, width, height = self.layout_manager.calculate_content_area(slide)
        
        current_y = top
        
        # Title
        if content.title:
            title_height = 1.0
            title_shape = slide.shapes.add_textbox(
                Inches(left), Inches(current_y), 
                Inches(width), Inches(title_height)
            )
            title_frame = title_shape.text_frame
            title_frame.text = content.title
            title_frame.paragraphs[0].font.size = Pt(self.config.typography.title_size)
            title_frame.paragraphs[0].font.name = self.config.typography.title_font
            title_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                self.config.color_scheme.primary.lstrip('#')
            )
            title_frame.paragraphs[0].font.bold = True
            current_y += title_height + 0.2
        
        # Subtitle
        if content.subtitle:
            subtitle_height = 0.6
            subtitle_shape = slide.shapes.add_textbox(
                Inches(left), Inches(current_y), 
                Inches(width), Inches(subtitle_height)
            )
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.text = content.subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(self.config.typography.subtitle_size)
            subtitle_frame.paragraphs[0].font.name = self.config.typography.body_font
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                self.config.color_scheme.text_secondary.lstrip('#')
            )
            current_y += subtitle_height + 0.2
        
        # Body content
        remaining_height = top + height - current_y
        self._add_body_content(slide, content.body, left, current_y, width, remaining_height)
    
    def _add_body_content(
        self, 
        slide: Slide, 
        body: List[Dict[str, Any]], 
        x: float, y: float, width: float, height: float
    ) -> None:
        """Add body content to slide."""
        current_y = y
        content_height_per_item = height / max(len(body), 1)
        
        for item in body:
            item_height = min(content_height_per_item, height - (current_y - y))
            if item_height <= 0:
                break
            
            content_type = item.get('type', 'text')
            
            if content_type == 'text':
                self._add_text_content(slide, item, x, current_y, width, item_height)
            elif content_type == 'bullet_list':
                self._add_bullet_list(slide, item, x, current_y, width, item_height)
            elif content_type == 'image':
                self._add_image_content(slide, item, x, current_y, width, item_height)
            elif content_type == 'table':
                self._add_table_content(slide, item, x, current_y, width, item_height)
            elif content_type == 'equation':
                self._add_equation_content(slide, item, x, current_y, width, item_height)
            elif content_type == 'chart':
                self._add_chart_content(slide, item, x, current_y, width, item_height)
            
            current_y += item_height
    
    def _add_text_content(
        self, slide: Slide, item: Dict[str, Any], 
        x: float, y: float, width: float, height: float
    ) -> None:
        """Add text content to slide."""
        text_content = item.get('content', '')
        
        # Process equations in text
        equations = self.equation_renderer.extract_equations_from_text(text_content)
        if equations:
            # Handle text with equations
            self._add_text_with_equations(slide, text_content, equations, x, y, width, height)
        else:
            # Regular text
            text_box = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            text_frame = text_box.text_frame
            text_frame.text = text_content
            text_frame.word_wrap = True
            
            # Apply formatting
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(self.config.typography.body_size)
                paragraph.font.name = self.config.typography.body_font
                paragraph.font.color.rgb = RGBColor.from_string(
                    self.config.color_scheme.text_primary.lstrip('#')
                )
    
    def _add_bullet_list(
        self, slide: Slide, item: Dict[str, Any], 
        x: float, y: float, width: float, height: float
    ) -> None:
        """Add bullet list to slide."""
        text_box = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.clear()
        
        bullets = item.get('items', [])
        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet_text
            p.level = item.get('level', 0)
            p.font.size = Pt(self.config.typography.body_size)
            p.font.name = self.config.typography.body_font
            p.font.color.rgb = RGBColor.from_string(
                self.config.color_scheme.text_primary.lstrip('#')
            )
    
    def _add_image_content(
        self, slide: Slide, item: Dict[str, Any], 
        x: float, y: float, width: float, height: float
    ) -> None:
        """Add image content to slide."""
        image_path = item.get('path') or item.get('url')
        if not image_path:
            return
        
        try:
            # Process image
            target_size = (int(width * 72), int(height * 72))  # Convert to pixels
            image_buffer = self.image_processor.process_image(image_path, target_size)
            
            # Add to slide
            pic = slide.shapes.add_picture(
                image_buffer, Inches(x), Inches(y), 
                width=Inches(width), height=Inches(height)
            )
            
            # Add caption if provided
            caption = item.get('caption')
            if caption:
                caption_y = y + height
                caption_height = 0.4
                caption_box = slide.shapes.add_textbox(
                    Inches(x), Inches(caption_y), 
                    Inches(width), Inches(caption_height)
                )
                caption_frame = caption_box.text_frame
                caption_frame.text = caption
                caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                caption_frame.paragraphs[0].font.size = Pt(self.config.typography.caption_size)
                caption_frame.paragraphs[0].font.name = self.config.typography.body_font
                caption_frame.paragraphs[0].font.italic = True
                
        except Exception as e:
            self.logger.error(f"Failed to add image: {e}")
            # Add placeholder
            placeholder_box = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            placeholder_box.text_frame.text = f"[Image: {item.get('alt_text', 'Image placeholder')}]"
    
    def _add_equation_content(
        self, slide: Slide, item: Dict[str, Any], 
        x: float, y: float, width: float, height: float
    ) -> None:
        """Add equation content to slide."""
        latex_code = item.get('latex', item.get('content', ''))
        
        try:
            # Render equation to image
            equation_image = self.equation_renderer.render_latex_to_image(
                latex_code, font_size=int(self.config.typography.body_size * 1.2)
            )
            
            if equation_image:
                # Add equation image
                slide.shapes.add_picture(
                    equation_image, Inches(x), Inches(y),
                    width=Inches(min(width, 6)), height=Inches(min(height, 2))
                )
            else:
                # Fallback to text
                text_box = slide.shapes.add_textbox(
                    Inches(x), Inches(y), Inches(width), Inches(height)
                )
                text_box.text_frame.text = f"Equation: {latex_code}"
                
        except Exception as e:
            self.logger.error(f"Failed to add equation: {e}")
            # Add placeholder
            text_box = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            text_box.text_frame.text = f"[Equation: {latex_code}]"
    
    def _add_chart_content(
        self, slide: Slide, item: Dict[str, Any], 
        x: float, y: float, width: float, height: float
    ) -> None:
        """Add chart content to slide."""
        chart_type = item.get('chart_type', 'column')
        chart_data = item.get('data', {})
        
        try:
            # Prepare chart data
            chart_data_obj = CategoryChartData()
            
            categories = chart_data.get('categories', [])
            series_data = chart_data.get('series', [])
            
            chart_data_obj.categories = categories
            for series in series_data:
                chart_data_obj.add_series(series.get('name', ''), series.get('values', []))
            
            # Determine chart type
            chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
            if chart_type.lower() == 'line':
                chart_type_enum = XL_CHART_TYPE.LINE
            elif chart_type.lower() == 'pie':
                chart_type_enum = XL_CHART_TYPE.PIE
            elif chart_type.lower() == 'bar':
                chart_type_enum = XL_CHART_TYPE.BAR_CLUSTERED
            
            # Add chart
            chart = slide.shapes.add_chart(
                chart_type_enum, Inches(x), Inches(y), 
                Inches(width), Inches(height), chart_data_obj
            ).chart
            
            # Customize chart appearance
            chart.has_legend = True
            chart.legend.position = 2  # Right side
            
        except Exception as e:
            self.logger.error(f"Failed to add chart: {e}")
            # Add placeholder
            text_box = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(width), Inches(height)
            )
            text_box.text_frame.text = f"[Chart: {item.get('title', 'Chart placeholder')}]"
    
    def _add_text_with_equations(
        self, slide: Slide, text: str, equations: List[Tuple[str, str]],
        x: float, y: float, width: float, height: float
    ) -> None:
        """Add text content that contains equations."""
        # For now, render as regular text with equation placeholders
        # In production, you might want to create multiple text boxes and image objects
        processed_text = text
        for eq_text, latex_code in equations:
            processed_text = processed_text.replace(eq_text, f"[Equation: {latex_code}]")
        
        text_box = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(width), Inches(height)
        )
        text_frame = text_box.text_frame
        text_frame.text = processed_text
        text_frame.word_wrap = True
    
    def _add_branding_elements(self, slide: Slide, slide_number: int) -> None:
        """Add branding elements to slide."""
        # Add logo
        if self.config.branding.logo_path:
            self._add_logo(slide)
        
        # Add slide number
        if self.config.branding.show_slide_numbers:
            self._add_slide_number(slide, slide_number)
        
        # Add custom footer
        if self.config.branding.custom_footer:
            self._add_footer(slide)
        
        # Add date
        if self.config.branding.show_date:
            self._add_date(slide)
    
    def _add_logo(self, slide: Slide) -> None:
        """Add logo to slide."""
        if not self.config.branding.logo_path:
            return
        
        try:
            logo_width, logo_height = self.config.branding.logo_size
            position = self.config.branding.logo_position
            
            # Calculate logo position
            if position == "top_left":
                x, y = 0.2, 0.2
            elif position == "top_right":
                x = self.config.slide_size[0] - logo_width - 0.2
                y = 0.2
            elif position == "bottom_left":
                x = 0.2
                y = self.config.slide_size[1] - logo_height - 0.2
            else:  # bottom_right
                x = self.config.slide_size[0] - logo_width - 0.2
                y = self.config.slide_size[1] - logo_height - 0.2
            
            # Process and add logo
            logo_buffer = self.image_processor.process_image(
                self.config.branding.logo_path,
                target_size=(int(logo_width * 72), int(logo_height * 72))
            )
            
            slide.shapes.add_picture(
                logo_buffer, Inches(x), Inches(y),
                width=Inches(logo_width), height=Inches(logo_height)
            )
            
        except Exception as e:
            self.logger.error(f"Failed to add logo: {e}")
    
    def _add_slide_number(self, slide: Slide, slide_number: int) -> None:
        """Add slide number to slide."""
        try:
            x = self.config.slide_size[0] - 1.0
            y = self.config.slide_size[1] - 0.5
            
            number_box = slide.shapes.add_textbox(
                Inches(x), Inches(y), Inches(0.8), Inches(0.3)
            )
            number_frame = number_box.text_frame
            number_frame.text = str(slide_number)
            number_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            number_frame.paragraphs[0].font.size = Pt(self.config.typography.footnote_size)
            number_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                self.config.color_scheme.text_secondary.lstrip('#')
            )
            
        except Exception as e:
            self.logger.error(f"Failed to add slide number: {e}")
    
    def _add_footer(self, slide: Slide) -> None:
        """Add custom footer to slide."""
        if not self.config.branding.custom_footer:
            return
        
        try:
            footer_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(self.config.slide_size[1] - 0.5),
                Inches(self.config.slide_size[0] - 1.0), Inches(0.3)
            )
            footer_frame = footer_box.text_frame
            footer_frame.text = self.config.branding.custom_footer
            footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            footer_frame.paragraphs[0].font.size = Pt(self.config.typography.footnote_size)
            footer_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                self.config.color_scheme.text_secondary.lstrip('#')
            )
            
        except Exception as e:
            self.logger.error(f"Failed to add footer: {e}")
    
    def _add_date(self, slide: Slide) -> None:
        """Add date to slide."""
        try:
            from datetime import datetime
            current_date = datetime.now().strftime("%B %d, %Y")
            
            date_box = slide.shapes.add_textbox(
                Inches(0.2), Inches(self.config.slide_size[1] - 0.5),
                Inches(2.0), Inches(0.3)
            )
            date_frame = date_box.text_frame
            date_frame.text = current_date
            date_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            date_frame.paragraphs[0].font.size = Pt(self.config.typography.footnote_size)
            date_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                self.config.color_scheme.text_secondary.lstrip('#')
            )
            
        except Exception as e:
            self.logger.error(f"Failed to add date: {e}")
    
    def _apply_slide_transition(self, slide: Slide) -> None:
        """Apply transition effects to slide."""
        # Note: python-pptx has limited transition support
        # Advanced transitions would require direct XML manipulation
        pass
    
    def _add_references_slide(self, citations: List[Citation]) -> None:
        """Add references slide to presentation."""
        if not self.presentation or not citations:
            return
        
        try:
            # Add references slide
            slide_layout = self.presentation.slide_layouts[6]  # Blank layout
            slide = self.presentation.slides.add_slide(slide_layout)
            
            left, top, width, height = self.layout_manager.calculate_content_area(slide)
            
            # Title
            title_shape = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(1.0)
            )
            title_frame = title_shape.text_frame
            title_frame.text = "References"
            title_frame.paragraphs[0].font.size = Pt(self.config.typography.title_size)
            title_frame.paragraphs[0].font.name = self.config.typography.title_font
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                self.config.color_scheme.primary.lstrip('#')
            )
            
            # References content
            refs_top = top + 1.2
            refs_height = height - 1.2
            
            refs_box = slide.shapes.add_textbox(
                Inches(left), Inches(refs_top), Inches(width), Inches(refs_height)
            )
            refs_frame = refs_box.text_frame
            refs_frame.clear()
            
            # Format citations
            citation_style = 'ieee'  # Could be configurable
            for i, citation in enumerate(citations):
                if i == 0:
                    p = refs_frame.paragraphs[0]
                else:
                    p = refs_frame.add_paragraph()
                
                formatted_citation = self.citation_formatter.format_citation(citation, citation_style)
                p.text = f"[{i+1}] {formatted_citation}"
                p.font.size = Pt(self.config.typography.body_size - 4)
                p.font.name = self.config.typography.body_font
                p.space_after = Pt(6)
            
            # Add branding
            slide_number = len(self.presentation.slides)
            self._add_branding_elements(slide, slide_number)
            
        except Exception as e:
            self.logger.error(f"Failed to add references slide: {e}")
    
    def _apply_global_formatting(self) -> None:
        """Apply global formatting to the presentation."""
        if not self.presentation:
            return
        
        try:
            # Set default fonts and colors for the presentation
            # This would typically involve modifying the slide master
            # For now, we rely on individual slide formatting
            pass
            
        except Exception as e:
            self.logger.error(f"Failed to apply global formatting: {e}")
    
    def _add_metadata(self, metadata: Dict[str, Any]) -> None:
        """Add metadata to presentation properties."""
        if not self.presentation:
            return
        
        try:
            core_props = self.presentation.core_properties
            
            if 'title' in metadata:
                core_props.title = metadata['title']
            if 'author' in metadata:
                core_props.author = metadata['author']
            if 'subject' in metadata:
                core_props.subject = metadata['subject']
            if 'description' in metadata:
                core_props.comments = metadata['description']
            if 'keywords' in metadata:
                core_props.keywords = metadata['keywords']
            if 'category' in metadata:
                core_props.category = metadata['category']
                
        except Exception as e:
            self.logger.error(f"Failed to add metadata: {e}")
    
    def _add_error_slide(self, error_message: str) -> None:
        """Add an error slide when content generation fails."""
        if not self.presentation:
            return
        
        try:
            slide_layout = self.presentation.slide_layouts[6]  # Blank layout
            slide = self.presentation.slides.add_slide(slide_layout)
            
            left, top, width, height = self.layout_manager.calculate_content_area(slide)
            
            error_box = slide.shapes.add_textbox(
                Inches(left), Inches(top + height/3), 
                Inches(width), Inches(height/3)
            )
            error_frame = error_box.text_frame
            error_frame.text = f"Error: {error_message}"
            error_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            error_frame.paragraphs[0].font.size = Pt(self.config.typography.body_size)
            error_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)  # Red text
            
        except Exception as e:
            self.logger.error(f"Failed to add error slide: {e}")
    
    def export_to_buffer(self, slides: List[SlideContent], **kwargs) -> io.BytesIO:
        """
        Export presentation to BytesIO buffer.
        
        Args:
            slides: List of slide content
            **kwargs: Additional parameters (citations, metadata, etc.)
            
        Returns:
            BytesIO buffer containing the PPTX file
        """
        citations = kwargs.get('citations')
        metadata = kwargs.get('metadata')
        
        # Create presentation
        presentation = self.create_presentation(slides, citations, metadata)
        
        # Save to buffer
        buffer = io.BytesIO()
        self.save_presentation(buffer)
        buffer.seek(0)
        
        return buffer
    
    def export_to_file(
        self, 
        slides: List[SlideContent], 
        output_path: str, 
        **kwargs
    ) -> None:
        """
        Export presentation to file.
        
        Args:
            slides: List of slide content
            output_path: Output file path
            **kwargs: Additional parameters (citations, metadata, etc.)
        """
        citations = kwargs.get('citations')
        metadata = kwargs.get('metadata')
        
        # Create presentation
        presentation = self.create_presentation(slides, citations, metadata)
        
        # Save to file
        self.save_presentation(output_path)


def create_academic_template_config(
    template: AcademicTemplate = AcademicTemplate.IEEE,
    university_name: Optional[str] = None,
    logo_path: Optional[str] = None,
    custom_colors: Optional[Dict[str, str]] = None
) -> TemplateConfig:
    """
    Create a template configuration for academic presentations.
    
    Args:
        template: Academic template to use
        university_name: University name for branding
        logo_path: Path to university logo
        custom_colors: Custom color overrides
        
    Returns:
        Configured TemplateConfig object
    """
    config = TemplateConfig()
    config.template = template
    
    # Apply university branding
    if university_name:
        config.branding.university_name = university_name
        config.branding.custom_footer = university_name
    
    if logo_path:
        config.branding.logo_path = logo_path
    
    # Apply custom colors
    if custom_colors:
        for key, value in custom_colors.items():
            if hasattr(config.color_scheme, key):
                setattr(config.color_scheme, key, value)
    
    return config


# Example usage and factory functions
def create_ieee_presentation(
    slides: List[SlideContent],
    citations: Optional[List[Citation]] = None,
    university_name: Optional[str] = None,
    logo_path: Optional[str] = None
) -> io.BytesIO:
    """Create an IEEE-style academic presentation."""
    config = create_academic_template_config(
        AcademicTemplate.IEEE, university_name, logo_path
    )
    generator = PPTXGenerator(config)
    return generator.export_to_buffer(slides, citations=citations)


def create_nature_presentation(
    slides: List[SlideContent],
    citations: Optional[List[Citation]] = None,
    university_name: Optional[str] = None,
    logo_path: Optional[str] = None
) -> io.BytesIO:
    """Create a Nature-style academic presentation."""
    config = create_academic_template_config(
        AcademicTemplate.NATURE, university_name, logo_path
    )
    generator = PPTXGenerator(config)
    return generator.export_to_buffer(slides, citations=citations)


def create_custom_presentation(
    slides: List[SlideContent],
    template_config: TemplateConfig,
    citations: Optional[List[Citation]] = None,
    metadata: Optional[Dict[str, Any]] = None
) -> io.BytesIO:
    """Create a presentation with custom template configuration."""
    generator = PPTXGenerator(template_config)
    return generator.export_to_buffer(
        slides, citations=citations, metadata=metadata
    )